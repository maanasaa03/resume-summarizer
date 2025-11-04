import os
import json
import re
import streamlit as st
from dotenv import load_dotenv
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from google import genai
import pdfplumber

load_dotenv()
client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))

# ---------- Enhanced text extraction with robust cleaning ----------
def extract_text(file_path):
    def clean_text(text):
        # Remove Word XML artifacts and special characters
        xml_artifacts = [
            r"_x[0-9A-Fa-f]{4}_",  # Word XML hex codes
            r"<[^>]+>",            # Any XML/HTML tags
            r"&\w+;",              # XML entities
            r"\u0000",             # Null characters
            r"\u0001-\u0008",      # Control characters
            r"\u000B-\u000C",      # More control characters
            r"\u000E-\u001F",      # Additional control characters
        ]
        
        for pattern in xml_artifacts:
            text = re.sub(pattern, "", text)
        
        # Replace various whitespace characters with single spaces
        whitespace_patterns = [
            r"\t+",                # Tabs
            r"\r\n",               # Windows line endings
            r"\r",                 # Old Mac line endings
            r"\n+",                # Multiple newlines
            r"\f",                 # Form feeds
            r"\v",                 # Vertical tabs
            r"\u200B",             # Zero-width space
            r"\u00A0",             # Non-breaking space
            r"\u2000-\u200F",      # Various space characters
            r"\u2028-\u2029",      # Line/paragraph separators
        ]
        
        for pattern in whitespace_patterns:
            text = re.sub(pattern, " ", text)
        
        # Remove page breaks and section breaks indicators
        page_break_indicators = [
            r"-\s*Page\s*\d+\s*-",  # Page indicators like "- Page 1 -"
            r"\[page \d+\]",        # [page 1] style
            r"\x0c",                # Form feed (page break)
            r"\f",                  # Form feed alternative
        ]
        
        for pattern in page_break_indicators:
            text = re.sub(pattern, " ", text, flags=re.IGNORECASE)
        
        # Clean up multiple spaces and trim
        text = re.sub(r" +", " ", text)
        text = text.strip()
        
        # Remove duplicate lines and excessive blank lines
        lines = text.split('\n')
        seen = set()
        unique_lines = []
        for line in lines:
            clean_line = line.strip()
            if clean_line and clean_line not in seen:
                seen.add(clean_line)
                unique_lines.append(clean_line)
        
        return '\n'.join(unique_lines)

    try:
        if file_path.lower().endswith(".pdf"):
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    # Try multiple extraction strategies for PDF
                    page_text = page.extract_text()
                    if not page_text or len(page_text.strip()) < 10:
                        # Fallback: extract tables or use layout preservation
                        page_text = page.extract_text(layout=True)
                    
                    if page_text:
                        # Clean up PDF-specific artifacts
                        page_text = re.sub(r'\s*-\s*\n\s*', '-', page_text)  # Handle hyphenated words
                        page_text = re.sub(r'\n\s*(?=[a-z])', ' ', page_text)  # Join broken lines
                        text += page_text + "\n"
            
            return clean_text(text)
            
        else:  # Word document
            doc = Document(file_path)
            full_text = []
            
            # Extract from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text)
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            full_text.append(cell.text)
            
            # Extract from headers and footers (if accessible)
            try:
                for section in doc.sections:
                    if section.header:
                        for paragraph in section.header.paragraphs:
                            if paragraph.text.strip():
                                full_text.append(paragraph.text)
                    if section.footer:
                        for paragraph in section.footer.paragraphs:
                            if paragraph.text.strip():
                                full_text.append(paragraph.text)
            except:
                pass  # Headers/footers might not be accessible in all documents
            
            text = "\n".join(full_text)
            return clean_text(text)
            
    except Exception as e:
        st.error(f"Error extracting text from file: {str(e)}")
        return ""

# ---------- Summarize using Gemini ----------
def summarize_with_gemini(text):
    # Pre-process text to remove any remaining artifacts
    text = re.sub(r'[^\x00-\x7F]+', ' ', text)  # Remove non-ASCII characters
    text = re.sub(r'\s+', ' ', text).strip()

    prompt = f"""
You are a professional technical résumé summarizer creating concise, impactful summaries for a 1-page PowerPoint résumé.

Produce JSON with these exact keys:
{{
  "name": "",
  "role": "",
  "location": "",
  "profile_overview": "",
  "professional_experience": "",
  "skills": "",
  "domain_experience": "",
  "education_and_certification": ""
}}

🧠 Guidelines:

1. PROFILE OVERVIEW:
   - 5–6 lines in a paragraph highlighting experience, domains, tools/technologies (**bold**), and achievements.
   - Include domain expertise, key tools, automation/workflow knowledge.
   - Example: "Accomplished Data Engineer with over 5 years of experience designing and delivering robust end-to-end ETL/ELT pipelines across Finance, Healthcare and Banking domains. Proficient in **dbt**, **Snowflake**, **Azure Data Factory (ADF)**, **SQL**. Expert in automating ETL with **ControlM**, **MFP**, and **ADF**, integrating data from **Oracle**, **SQL Server**, and **AWS S3**."

2. LOCATION:
   - Extract candidate's current city and country. Always fill this field.
   - Example: "Bangalore, India"

3. PROFESSIONAL EXPERIENCE:
   - Include **only the 3 most recent projects or roles** from the résumé.
   - Each project should have a **bold heading** with 2–3 bullets explaining what was built, tools used, and outcomes.
   - Bullet points should start with "* ".
   - Limit total text to ~900–1000 characters to fit in PowerPoint.
   - Example:
     **Customer Data Pipeline Migration at XYZ Corp**
     * Leveraged **Snowflake**, **dbt**, **Azure Data Factory (ADF)** to build scalable data models and ensure high-quality reconciled datasets.
     * Automated ETL processes with **ControlM**, improving data processing efficiency by 30%.

4. SKILLS: Comma-separated list of key skills/technologies. Limit to 15 top skills.

5. DOMAIN EXPERIENCE: 2–3 short items describing domains worked in.

6. EDUCATION & CERTIFICATION: 1–2 lines. Always fill this field.

🔹 Formatting rules:
- Use Markdown-style bold (**text**) for tools, technologies, and project names.
- Keep text compact (~1000 characters per large section).
- Preserve line breaks for bullets and paragraphs.

Resume text:
{text}
"""
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt
    )

    content = response.text.strip()
    try:
        data = json.loads(content)
        for k, v in data.items():
            if isinstance(v, str):
                data[k] = v.replace('\\n', '\n').strip()
    except Exception:
        data = {}
        for match in re.finditer(r'"(\w+)"\s*:\s*"([^"]*)"', content, re.DOTALL):
            data[match.group(1)] = match.group(2).replace("\\n", "\n").strip()
    return data

# ---------- Apply bold for Markdown syntax ----------
def apply_bold_markdown(paragraph, text):
    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        else:
            run.text = part

# ---------- Auto-fit text ----------
def auto_fit_text(shape, max_size=11, min_size=9):
    tf = shape.text_frame
    total_chars = sum(len(p.text) for p in tf.paragraphs)
    shrink_factor = 1.0

    if total_chars > 800:
        shrink_factor = 0.9
    if total_chars > 1200:
        shrink_factor = 0.8
    if total_chars > 1600:
        shrink_factor = 0.7

    font_size = max(min_size, int(max_size * shrink_factor))
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)

# ---------- Populate PowerPoint Template ----------
def fill_ppt_template(template_path, data, output_path):
    from pptx.enum.text import MSO_ANCHOR

    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for key, val in data.items():
                placeholder = f"{{{{{key.upper()}}}}}"
                if placeholder in shape.text:
                    new_text = shape.text.replace(placeholder, val)
                    shape.text = ""
                    tf = shape.text_frame
                    tf.clear()
                    tf.word_wrap = True
                    tf.vertical_anchor = MSO_ANCHOR.TOP

                    for line in new_text.split("\n"):
                        p = tf.add_paragraph()
                        apply_bold_markdown(p, line)
                        p.space_after = Pt(4)

                    for paragraph in tf.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = "Aptos"
                            if key in ["name", "role", "location", "education_and_certification"]:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                            else:
                                run.font.color.rgb = RGBColor(40, 40, 40)

                    auto_fit_text(shape)

    prs.save(output_path)

# ---------- Streamlit UI ----------
st.title("📄 AI Resume to PowerPoint Generator")
st.write("Upload your resume (.pdf or .docx), then click the button to generate a summarized PowerPoint résumé using Gemini AI.")

uploaded_file = st.file_uploader("Choose a resume file", type=["pdf", "docx"])

if uploaded_file:
    process_button = st.button("🚀 Upload & Generate PowerPoint")

    if process_button:
        template_path = "refined_template.pptx"
        output_folder = "output"
        os.makedirs(output_folder, exist_ok=True)

        input_path = os.path.join(output_folder, uploaded_file.name)
        with open(input_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        with st.spinner("Extracting and summarizing your resume..."):
            text = extract_text(input_path)
            data = summarize_with_gemini(text)
            output_path = os.path.join(output_folder, os.path.splitext(uploaded_file.name)[0] + ".pptx")
            fill_ppt_template(template_path, data, output_path)

        st.success("✅ PowerPoint generated successfully!")
        with open(output_path, "rb") as f:
            st.download_button("📥 Download PowerPoint", f, file_name=os.path.basename(output_path))